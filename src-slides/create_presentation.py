#!/usr/bin/env python3
"""
PowerPoint Generator for .NET Agentic Modernization Demo Session

This script creates a PowerPoint presentation from the slide content and speaker notes
defined in the docs/05_slide-content-and-speaker-notes.md file.

Usage:
    python create_presentation.py [--template TEMPLATE_PATH] [--output OUTPUT_PATH]

Options:
    --template  Path to an existing PowerPoint file to use as a style template
    --output    Output path for the generated presentation (default: presentation.pptx)

Examples:
    # Generate presentation with default styling
    python create_presentation.py
    
    # Generate presentation using an existing template for styling
    python create_presentation.py --template company_template.pptx
    
    # Specify custom output path
    python create_presentation.py --output my_presentation.pptx
"""

import argparse
import os
import re
import sys
from datetime import datetime
from pathlib import Path
from typing import List, Optional, Tuple

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Error: python-pptx is required. Install it with: pip install python-pptx")
    sys.exit(1)


def parse_slide_content(markdown_path: str) -> List[dict]:
    """
    Parse the markdown file to extract slide content.
    
    Returns a list of dictionaries with:
    - slide_number: int
    - title: str
    - content: list of str (bullet points)
    - speaker_notes: str
    - code_block: Optional[str]
    """
    slides = []
    
    with open(markdown_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Split by slide headers (## Slide N: Title)
    slide_pattern = r'## Slide (\d+): (.+?)(?=\n## Slide \d+:|## Appendix|# End of Document|$)'
    matches = re.findall(slide_pattern, content, re.DOTALL)
    
    for match in matches:
        slide_num = int(match[0])
        slide_content = match[1]
        
        # Extract title from the first line
        lines = slide_content.strip().split('\n')
        title_line = lines[0] if lines else ""
        
        # Extract slide title from ### Slide Content section
        title_match = re.search(r'\*\*Title\*\*:\s*(.+)', slide_content)
        title = title_match.group(1).strip() if title_match else f"Slide {slide_num}"
        
        # Extract key points / bullet content
        key_points = []
        key_points_match = re.search(r'\*\*Key Points\*\*:\s*\n((?:- .+\n?)+)', slide_content)
        if key_points_match:
            key_points = [line.strip('- ').strip() for line in key_points_match.group(1).strip().split('\n') if line.strip()]
        
        # Extract code blocks
        code_block = None
        code_match = re.search(r'```(?:csharp|json)?\n(.+?)```', slide_content, re.DOTALL)
        if code_match:
            code_block = code_match.group(1).strip()
        
        # Extract speaker notes
        speaker_notes = ""
        notes_match = re.search(r'### Speaker Notes\n>\s*(.+?)(?=\n---|\n##|$)', slide_content, re.DOTALL)
        if notes_match:
            speaker_notes = notes_match.group(1).strip()
        
        # Extract subtitle if present
        subtitle = ""
        subtitle_match = re.search(r'\*\*Subtitle\*\*:\s*(.+)', slide_content)
        if subtitle_match:
            subtitle = subtitle_match.group(1).strip()
        
        # Extract definition if present
        definition = ""
        def_match = re.search(r'\*\*Definition\*\*:\s*(.+)', slide_content)
        if def_match:
            definition = def_match.group(1).strip()
        
        # Extract key characteristics
        characteristics = []
        char_match = re.search(r'\*\*Key Characteristics\*\*:\s*\n((?:- .+\n?)+)', slide_content)
        if char_match:
            characteristics = [line.strip('- ').strip() for line in char_match.group(1).strip().split('\n') if line.strip()]
        
        # Extract benefits
        benefits = []
        benefits_match = re.search(r'\*\*Benefits\*\*:\s*\n((?:- .+\n?)+)', slide_content)
        if benefits_match:
            benefits = [line.strip('- ').strip() for line in benefits_match.group(1).strip().split('\n') if line.strip()]
        
        # Extract "What We'll Do" items
        what_well_do = []
        wwd_match = re.search(r'\*\*What We\'ll Do\*\*:\s*\n((?:\d+\. .+\n?)+)', slide_content)
        if wwd_match:
            what_well_do = [re.sub(r'^\d+\.\s*', '', line).strip() for line in wwd_match.group(1).strip().split('\n') if line.strip()]
        
        # Extract features
        features = []
        features_match = re.search(r'\*\*Features\*\*:\s*\n((?:- .+\n?)+)', slide_content)
        if features_match:
            features = [line.strip('- ').strip() for line in features_match.group(1).strip().split('\n') if line.strip()]
        
        # Extract "What's Missing"
        whats_missing = []
        missing_match = re.search(r"\*\*What's Missing\*\*:\s*\n((?:- .+\n?)+)", slide_content)
        if missing_match:
            whats_missing = [line.strip('- ').strip() for line in missing_match.group(1).strip().split('\n') if line.strip()]
        
        # Extract Core Technologies
        core_tech = []
        tech_match = re.search(r'\*\*Core Technologies\*\*:\s*\n((?:- .+\n?)+)', slide_content)
        if tech_match:
            core_tech = [line.strip('- ').strip() for line in tech_match.group(1).strip().split('\n') if line.strip()]
        
        # Extract Workflow items
        workflow = []
        workflow_match = re.search(r'\*\*Workflow\*\*:\s*\n((?:\d+\. .+\n?)+)', slide_content)
        if workflow_match:
            workflow = [re.sub(r'^\d+\.\s*', '', line).strip() for line in workflow_match.group(1).strip().split('\n') if line.strip()]
        
        # Extract Azure AI Foundry Benefits
        foundry_benefits = []
        foundry_match = re.search(r'\*\*Azure AI Foundry Benefits\*\*:\s*\n((?:- .+\n?)+)', slide_content)
        if foundry_match:
            foundry_benefits = [line.strip('- ').strip() for line in foundry_match.group(1).strip().split('\n') if line.strip()]
        
        # Extract Best Practices
        best_practices = []
        bp_match = re.search(r'\*\*Best Practices\*\*:\s*\n((?:- .+\n?)+)', slide_content)
        if bp_match:
            best_practices = [line.strip('- ').strip() for line in bp_match.group(1).strip().split('\n') if line.strip()]
        
        # Combine all bullet points
        all_bullets = key_points + characteristics + benefits + what_well_do + features + whats_missing + core_tech + workflow + foundry_benefits + best_practices
        
        slides.append({
            'slide_number': slide_num,
            'title': title,
            'subtitle': subtitle,
            'definition': definition,
            'content': all_bullets,
            'code_block': code_block,
            'speaker_notes': speaker_notes
        })
    
    return slides


def create_presentation(slides: List[dict], template_path: Optional[str] = None, output_path: str = "presentation.pptx"):
    """
    Create a PowerPoint presentation from parsed slide data.
    
    Args:
        slides: List of slide dictionaries from parse_slide_content
        template_path: Optional path to a template PPTX file for styling
        output_path: Output path for the generated presentation
    """
    # Create presentation from template or new
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
        # Clear existing slides if using template
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
        print(f"Using template: {template_path}")
    else:
        prs = Presentation()
        if template_path:
            print(f"Template not found: {template_path}. Using default styling.")
    
    # Set slide dimensions (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for slide_data in slides:
        # Choose layout based on slide type
        if slide_data['slide_number'] == 1:
            # Title slide
            layout = prs.slide_layouts[6]  # Blank layout for custom title
            slide = prs.slides.add_slide(layout)
            
            # Add title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1.5))
            title_frame = title_box.text_frame
            title_frame.paragraphs[0].text = slide_data['title']
            title_frame.paragraphs[0].font.size = Pt(44)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Add subtitle
            if slide_data['subtitle']:
                sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(1))
                sub_frame = sub_box.text_frame
                sub_frame.paragraphs[0].text = slide_data['subtitle']
                sub_frame.paragraphs[0].font.size = Pt(24)
                sub_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
        elif slide_data['slide_number'] == 14:
            # Q&A / Closing slide
            layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(layout)
            
            # Add title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
            title_frame = title_box.text_frame
            title_frame.paragraphs[0].text = slide_data['title']
            title_frame.paragraphs[0].font.size = Pt(48)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Add "Thank You!" text
            thanks_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
            thanks_frame = thanks_box.text_frame
            thanks_frame.paragraphs[0].text = "Thank You!"
            thanks_frame.paragraphs[0].font.size = Pt(32)
            thanks_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
        elif slide_data['code_block']:
            # Slide with code block
            layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(layout)
            
            # Add title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
            title_frame = title_box.text_frame
            title_frame.paragraphs[0].text = slide_data['title']
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True
            
            # Add code block
            code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5))
            code_frame = code_box.text_frame
            code_frame.word_wrap = True
            
            # Style code as monospace
            p = code_frame.paragraphs[0]
            p.text = slide_data['code_block']
            p.font.name = 'Consolas'
            p.font.size = Pt(12)
            
            # Add a background rectangle for code
            code_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.4), Inches(1.2),
                Inches(12.5), Inches(5.2)
            )
            code_shape.fill.solid()
            code_shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
            code_shape.line.color.rgb = RGBColor(200, 200, 200)
            
            # Move code box to front
            spTree = slide.shapes._spTree
            code_box_element = code_box._element
            spTree.remove(code_box_element)
            spTree.append(code_box_element)
            
        else:
            # Standard content slide
            layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(layout)
            
            # Add title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
            title_frame = title_box.text_frame
            title_frame.paragraphs[0].text = slide_data['title']
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True
            
            # Add definition if present
            current_y = 1.2
            if slide_data['definition']:
                def_box = slide.shapes.add_textbox(Inches(0.5), Inches(current_y), Inches(12.333), Inches(0.6))
                def_frame = def_box.text_frame
                def_frame.paragraphs[0].text = slide_data['definition']
                def_frame.paragraphs[0].font.size = Pt(18)
                def_frame.paragraphs[0].font.italic = True
                current_y += 0.8
            
            # Add bullet points
            if slide_data['content']:
                content_box = slide.shapes.add_textbox(Inches(0.5), Inches(current_y), Inches(12.333), Inches(5.5))
                content_frame = content_box.text_frame
                content_frame.word_wrap = True
                
                for i, bullet in enumerate(slide_data['content']):
                    if i == 0:
                        p = content_frame.paragraphs[0]
                    else:
                        p = content_frame.add_paragraph()
                    
                    # Handle bold markers
                    clean_bullet = bullet.replace('**', '')
                    p.text = f"• {clean_bullet}"
                    p.font.size = Pt(18)
                    p.space_after = Pt(12)
        
        # Add speaker notes
        if slide_data['speaker_notes']:
            notes_slide = slide.notes_slide
            notes_frame = notes_slide.notes_text_frame
            notes_frame.text = slide_data['speaker_notes']
    
    # Save presentation
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")


def main():
    parser = argparse.ArgumentParser(
        description='Create PowerPoint presentation from markdown slide content',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__
    )
    parser.add_argument(
        '--template',
        type=str,
        help='Path to an existing PowerPoint file to use as a style template'
    )
    parser.add_argument(
        '--output',
        type=str,
        default=None,
        help='Output path for the generated presentation (default: presentation-yymmdd-hhmmss.pptx)'
    )
    parser.add_argument(
        '--markdown',
        type=str,
        default=None,
        help='Path to the markdown file with slide content (default: ../docs/05_slide-content-and-speaker-notes.md)'
    )
    
    args = parser.parse_args()
    
    # Determine output file path
    if args.output:
        output_path = args.output
    else:
        # Generate timestamped filename in format: presentation-yymmdd-hhmmss.pptx
        now = datetime.now()
        timestamp = now.strftime("%y%m%d-%H%M%S")
        output_path = f"presentation-{timestamp}.pptx"
    
    # Determine markdown file path
    if args.markdown:
        markdown_path = args.markdown
    else:
        # Default to the docs folder relative to this script
        script_dir = Path(__file__).parent
        markdown_path = script_dir.parent / 'docs' / '05_slide-content-and-speaker-notes.md'
    
    if not os.path.exists(markdown_path):
        print(f"Error: Markdown file not found: {markdown_path}")
        sys.exit(1)
    
    print(f"Parsing slide content from: {markdown_path}")
    slides = parse_slide_content(str(markdown_path))
    print(f"Found {len(slides)} slides")
    
    for slide in slides:
        print(f"  Slide {slide['slide_number']}: {slide['title']}")
    
    create_presentation(slides, args.template, output_path)


if __name__ == '__main__':
    main()
